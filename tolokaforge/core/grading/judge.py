"""LLM Judge for grading — supports both single-call and agentic (tool-using) modes.

Agentic mode gives the judge workspace tools to inspect the agent's output.
The judge persona and rubric come from the pack's grading.yaml — the harness
supplies only the universal mechanism.

Architecture:
  - Core tools (always available): list_files, read_file, glob_files,
    grep_workspace, run_shell, submit_grade
  - Tool packs (opt-in via grading.yaml ``tool_packs``):
    * "office": read_xlsx_cell, read_xlsx_range, read_docx_content,
      read_pptx_content, render_pptx_slides, render_pdf_pages
"""

import fnmatch
import hashlib
import json
import logging
import re
import shutil
import subprocess
from pathlib import Path
from typing import Any

from jsonschema import validate

from tolokaforge.core.model_client import LLMClient
from tolokaforge.core.models import Message, ModelConfig

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Path sandboxing
# ---------------------------------------------------------------------------


def _safe_resolve(workspace: Path, rel_path: str) -> Path | None:
    """Resolve a relative path within workspace. Returns None if outside."""
    target = (workspace / rel_path).resolve()
    if not str(target).startswith(str(workspace.resolve())):
        return None
    return target


# ===================================================================
# CORE TOOLS — always available to the agentic judge
# ===================================================================

# -- Tool definitions (OpenAI function-calling format) --

_CORE_TOOL_DEFS: list[dict[str, Any]] = [
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List files in a workspace directory.",
            "parameters": {
                "type": "object",
                "properties": {
                    "directory": {
                        "type": "string",
                        "description": "Relative directory path (default: workspace root).",
                    }
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": (
                "Read a file from the workspace. Returns text content. "
                "For XLSX files returns CSV of each sheet. "
                "For DOCX files returns paragraph text."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path relative to workspace."},
                    "max_lines": {
                        "type": "integer",
                        "description": "Maximum lines to return (default 500).",
                    },
                    "offset": {
                        "type": "integer",
                        "description": "1-based line to start from (default 1).",
                    },
                    "limit": {
                        "type": "integer",
                        "description": "Lines to return (0 = use max_lines).",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "glob_files",
            "description": "Find files matching a glob pattern. Fast file discovery.",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Glob pattern (e.g. '**/*.py', '*.xlsx').",
                    },
                    "path": {
                        "type": "string",
                        "description": "Subdirectory to search in (default: workspace root).",
                    },
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "grep_workspace",
            "description": (
                "Search workspace files by content regex and/or filename glob. "
                "Supports context lines, output modes, and case-insensitive matching."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Regex pattern to search in file contents.",
                    },
                    "file_glob": {
                        "type": "string",
                        "description": "Glob pattern to filter filenames (e.g. '*.py').",
                    },
                    "path": {
                        "type": "string",
                        "description": "Subdirectory to search in (default: workspace root).",
                    },
                    "context_lines": {
                        "type": "integer",
                        "description": "Lines of context before and after each match.",
                    },
                    "output_mode": {
                        "type": "string",
                        "description": "'content' (matching lines), 'files' (paths only), 'count'.",
                    },
                    "case_insensitive": {
                        "type": "boolean",
                        "description": "Case-insensitive matching.",
                    },
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_shell",
            "description": "Run a shell command in the workspace directory.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {
                        "type": "string",
                        "description": "Shell command to execute.",
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "Max seconds before killing the process (default 60).",
                    },
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "submit_grade",
            "description": "Submit your final grade. Call this when you have finished evaluating.",
            "parameters": {
                "type": "object",
                "properties": {
                    "score": {
                        "type": "number",
                        "description": "Score between 0.0 and 1.0.",
                    },
                    "reasoning": {
                        "type": "string",
                        "description": "Explanation of the grade with specific evidence.",
                    },
                },
                "required": ["score", "reasoning"],
            },
        },
    },
]


# -- Core tool executors --


def _exec_list_files(workspace: Path, args: dict) -> str:
    directory = args.get("directory", ".")
    target = _safe_resolve(workspace, directory)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists() or not target.is_dir():
        return f"Directory not found: {directory}"
    entries = []
    for p in sorted(target.iterdir()):
        if p.name.startswith("."):
            continue
        kind = "dir" if p.is_dir() else f"{p.stat().st_size}B"
        entries.append(f"  {p.name}  [{kind}]")
    if not entries:
        return f"{directory}/ (empty)"
    return f"{directory}/\n" + "\n".join(entries)


def _exec_read_file(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists():
        return f"File not found: {path}"

    max_lines = int(args.get("max_lines", 500))
    offset = int(args.get("offset", 1))
    limit = int(args.get("limit", 0))
    suffix = target.suffix.lower()

    # XLSX → CSV per sheet
    if suffix == ".xlsx":
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(target), data_only=True)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"=== Sheet: {sheet_name} ===")
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= max_lines:
                        lines.append(f"... (truncated at {max_lines} rows)")
                        break
                    lines.append(",".join(str(c) if c is not None else "" for c in row))
            wb.close()
            return "\n".join(lines)
        except ImportError:
            return f"[XLSX file at {path} — openpyxl not installed]"

    # DOCX → paragraph text
    if suffix == ".docx":
        try:
            from docx import Document

            doc = Document(str(target))
            return "\n".join(p.text for p in doc.paragraphs[:max_lines])
        except ImportError:
            return f"[DOCX file at {path} — python-docx not installed]"

    # Other binary
    if suffix in (".pdf", ".pptx", ".png", ".jpg", ".jpeg", ".gif", ".zip"):
        return f"[Binary file: {path} ({target.stat().st_size} bytes)]"

    # Text files — with offset/limit support
    try:
        content = target.read_text(encoding="utf-8", errors="replace")
    except UnicodeDecodeError:
        return f"Binary file: {path}"
    all_lines = content.splitlines(keepends=True)
    total = len(all_lines)
    start = max(0, offset - 1)
    cap = limit if limit > 0 else max_lines
    selected = all_lines[start : start + cap]
    result = "".join(selected)
    if (start + cap) < total:
        result += f"\n... ({total} total lines)"
    return result


def _exec_glob_files(workspace: Path, args: dict) -> str:
    pattern = args.get("pattern", "")
    subdir = args.get("path", ".")
    target = _safe_resolve(workspace, subdir)
    if target is None:
        return "Error: path outside workspace"
    if not target.is_dir():
        return f"Directory not found: {subdir}"
    matches: list[tuple[float, str]] = []
    base = workspace.resolve()
    for fp in target.glob(pattern):
        if not fp.is_file() or fp.name.startswith("."):
            continue
        try:
            rel = str(fp.relative_to(base))
        except ValueError:
            continue
        try:
            mtime = fp.stat().st_mtime
        except OSError:
            mtime = 0
        matches.append((mtime, rel))
    if not matches:
        return f"(no files matching '{pattern}')"
    matches.sort(key=lambda x: x[0], reverse=True)
    paths = [m[1] for m in matches[:500]]
    result = "\n".join(paths)
    if len(matches) > 500:
        result += f"\n... ({len(matches) - 500} more files)"
    return result


def _exec_grep_workspace(workspace: Path, args: dict) -> str:
    pattern = args.get("pattern", "")
    file_glob = args.get("file_glob", "")
    subdir = args.get("path", ".")
    context_lines = int(args.get("context_lines", 0))
    output_mode = args.get("output_mode", "content")
    case_insensitive = bool(args.get("case_insensitive", False))
    max_matches = 500

    if not pattern and not file_glob:
        return "Provide at least one of 'pattern' or 'file_glob'"

    search_root = _safe_resolve(workspace, subdir)
    if search_root is None:
        return "Error: path outside workspace"
    if not search_root.is_dir():
        return f"Search path not found: {subdir}"

    flags = re.IGNORECASE if case_insensitive else 0
    regex = re.compile(pattern, flags) if pattern else None
    base = workspace.resolve()
    matches: list[str] = []
    file_counts: dict[str, int] = {}
    total_matches = 0

    for fp in sorted(search_root.rglob("*")):
        if not fp.is_file() or fp.name.startswith("."):
            continue
        if any(p in (".git", "__pycache__", "node_modules") for p in fp.parts):
            continue
        try:
            rel = str(fp.relative_to(base))
        except ValueError:
            continue

        if file_glob and not fnmatch.fnmatch(fp.name, file_glob):
            continue

        if not regex:
            matches.append(rel)
            total_matches += 1
            if total_matches >= max_matches:
                break
            continue

        try:
            text = fp.read_text(encoding="utf-8", errors="replace")
        except Exception:
            continue

        lines = text.splitlines()
        file_match_count = 0

        if output_mode == "files":
            if regex.search(text):
                matches.append(rel)
                total_matches += 1
            if total_matches >= max_matches:
                break
            continue

        for lineno, line in enumerate(lines):
            if regex.search(line):
                file_match_count += 1
                if output_mode == "count":
                    continue
                if context_lines > 0:
                    start = max(0, lineno - context_lines)
                    end = min(len(lines), lineno + context_lines + 1)
                    for ctx_i in range(start, end):
                        marker = ">" if ctx_i == lineno else " "
                        matches.append(f"{rel}:{ctx_i + 1}:{marker} {lines[ctx_i].rstrip()}")
                    matches.append("--")
                else:
                    matches.append(f"{rel}:{lineno + 1}:{line.rstrip()}")
                total_matches += 1
                if total_matches >= max_matches:
                    break

        if file_match_count > 0 and output_mode == "count":
            file_counts[rel] = file_match_count
            total_matches += 1
        if total_matches >= max_matches:
            break

    if output_mode == "count":
        if not file_counts:
            return "(no matches)"
        return "\n".join(f"{fp}: {count}" for fp, count in sorted(file_counts.items()))
    if not matches:
        return "(no matches)"
    result = "\n".join(matches)
    if total_matches >= max_matches:
        result += f"\n... (truncated at {max_matches} matches)"
    return result


def _exec_run_shell(workspace: Path, args: dict) -> str:
    command = args.get("command", "")
    timeout = int(args.get("timeout", 60))
    if not command:
        return "Error: no command provided"
    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=str(workspace),
        )
        output = result.stdout + result.stderr
        if len(output) > 50_000:
            output = output[:50_000] + "\n... (output truncated at 50000 chars)"
        return output or "(no output)"
    except subprocess.TimeoutExpired:
        return f"ERROR: Command timed out ({timeout}s limit)"


_CORE_EXECUTORS: dict[str, Any] = {
    "list_files": _exec_list_files,
    "read_file": _exec_read_file,
    "glob_files": _exec_glob_files,
    "grep_workspace": _exec_grep_workspace,
    "run_shell": _exec_run_shell,
}


# ===================================================================
# TOOL PACKS — opt-in via grading.yaml tool_packs field
# ===================================================================

# -- Office tool pack --

_OFFICE_TOOL_DEFS: list[dict[str, Any]] = [
    {
        "type": "function",
        "function": {
            "name": "read_xlsx_cell",
            "description": (
                "Read a single cell value and formula from an XLSX file. "
                "Returns both the computed value and the formula if present."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "XLSX file path relative to workspace.",
                    },
                    "sheet": {"type": "string", "description": "Sheet name."},
                    "cell": {
                        "type": "string",
                        "description": "Cell reference (e.g. 'B5', 'AA12').",
                    },
                },
                "required": ["path", "sheet", "cell"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_xlsx_range",
            "description": "Read a range of cells from an XLSX file as CSV.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "XLSX file path relative to workspace.",
                    },
                    "sheet": {"type": "string", "description": "Sheet name."},
                    "range_ref": {
                        "type": "string",
                        "description": "Cell range e.g. 'A1:D10'.",
                    },
                },
                "required": ["path", "sheet", "range_ref"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_docx_content",
            "description": (
                "Convert a DOCX file to markdown. Preserves heading hierarchy, "
                "tables, lists, and formatting. Uses pandoc if available, else python-docx."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "DOCX file path relative to workspace.",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Maximum characters to return (default 50000).",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_pptx_content",
            "description": (
                "Read structured text from a PPTX file: slide titles, body text, "
                "tables, chart data, and speaker notes."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "PPTX file path relative to workspace.",
                    },
                    "max_slides": {
                        "type": "integer",
                        "description": "Maximum slides to read (default 50).",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "render_pptx_slides",
            "description": (
                "Render PPTX slides to PNG images and extract text. "
                "Returns JSON with per-slide image_path and text."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "PPTX file path relative to workspace.",
                    },
                    "max_slides": {
                        "type": "integer",
                        "description": "Maximum slides to render (default 10).",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "render_pdf_pages",
            "description": (
                "Render PDF pages to PNG images and extract text. "
                "Returns JSON with per-page image_path and text."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "PDF file path relative to workspace.",
                    },
                    "max_pages": {
                        "type": "integer",
                        "description": "Maximum pages to render (default 10).",
                    },
                },
                "required": ["path"],
            },
        },
    },
]


# -- Office tool executors --


def _exec_read_xlsx_cell(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists():
        return f"File not found: {path}"
    sheet_name = args.get("sheet", "")
    cell_ref = args.get("cell", "")
    try:
        import openpyxl
    except ImportError:
        return "openpyxl not installed — install with: pip install tolokaforge[office]"
    try:
        wb_formula = openpyxl.load_workbook(str(target), data_only=False)
        if sheet_name not in wb_formula.sheetnames:
            avail = wb_formula.sheetnames
            wb_formula.close()
            return f"Sheet '{sheet_name}' not found. Available: {avail}"
        formula_val = wb_formula[sheet_name][cell_ref].value
        wb_formula.close()

        wb_data = openpyxl.load_workbook(str(target), data_only=True)
        data_val = wb_data[sheet_name][cell_ref].value
        wb_data.close()

        parts = [f"{path}!{sheet_name}!{cell_ref}:"]
        parts.append(f"  value: {data_val!r}")
        if isinstance(formula_val, str) and formula_val.startswith("="):
            parts.append(f"  formula: {formula_val}")
        return "\n".join(parts)
    except Exception as e:
        return f"Error reading {path}: {e}"


def _exec_read_xlsx_range(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists():
        return f"File not found: {path}"
    sheet_name = args.get("sheet", "")
    range_ref = args.get("range_ref", "")
    try:
        import openpyxl
    except ImportError:
        return "openpyxl not installed — install with: pip install tolokaforge[office]"
    try:
        wb = openpyxl.load_workbook(str(target), data_only=True)
        if sheet_name not in wb.sheetnames:
            avail = wb.sheetnames
            wb.close()
            return f"Sheet '{sheet_name}' not found. Available: {avail}"
        ws = wb[sheet_name]
        lines = []
        for row in ws[range_ref]:
            line = ",".join(str(c.value) if c.value is not None else "" for c in row)
            lines.append(line)
        wb.close()
        if len(lines) > 500:
            lines = lines[:500]
            lines.append("... (truncated at 500 rows)")
        return "\n".join(lines)
    except Exception as e:
        return f"Error reading {path}: {e}"


def _exec_read_docx_content(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists():
        return f"File not found: {path}"
    max_chars = int(args.get("max_chars", 50_000))

    # Try pandoc first
    pandoc = shutil.which("pandoc")
    if pandoc:
        try:
            result = subprocess.run(
                [pandoc, str(target), "-t", "markdown", "--wrap=none"],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and result.stdout.strip():
                content = result.stdout
                if len(content) > max_chars:
                    content = content[:max_chars] + "\n\n... (truncated)"
                return content
        except (subprocess.TimeoutExpired, Exception):
            pass

    # Fallback: python-docx
    try:
        from docx import Document
    except ImportError:
        return "Neither pandoc nor python-docx available — install with: pip install tolokaforge[office]"
    try:
        doc = Document(str(target))
        lines: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)
        for table in doc.tables:
            lines.append("")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
        content = "\n".join(lines)
        if len(content) > max_chars:
            content = content[:max_chars] + "\n\n... (truncated)"
        return content
    except Exception as e:
        return f"Error reading {path}: {e}"


def _exec_read_pptx_content(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return "Error: path outside workspace"
    if not target.exists():
        return f"File not found: {path}"
    max_slides = int(args.get("max_slides", 50))
    try:
        from pptx import Presentation
    except ImportError:
        return "python-pptx not installed — install with: pip install tolokaforge[office]"
    try:
        prs = Presentation(str(target))
        sections: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            if i > max_slides:
                break
            parts: list[str] = [f"=== Slide {i} ==="]
            if slide.shapes.title and slide.shapes.title.text:
                parts.append(f"Title: {slide.shapes.title.text}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != (slide.shapes.title.text if slide.shapes.title else ""):
                            parts.append(text)
                if shape.has_table:
                    table = shape.table
                    parts.append("[Table]")
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        parts.append(f"  {row_text}")
                if shape.has_chart:
                    try:
                        chart = shape.chart
                        parts.append(f"[Chart: {chart.chart_type}]")
                        for series in chart.series:
                            if hasattr(series, "values") and series.values:
                                parts.append(f"  Series values: {list(series.values)[:20]}")
                    except Exception:
                        parts.append("[Chart present — data extraction failed]")
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes] {notes}")
            sections.append("\n".join(parts))
        return "\n\n".join(sections)[:50_000]
    except Exception as e:
        return f"Error reading {path}: {e}"


def _exec_render_pptx_slides(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return json.dumps({"error": "path outside workspace"})
    if not target.exists():
        return json.dumps({"error": f"File not found: {path}"})
    max_slides = int(args.get("max_slides", 10))

    try:
        from pptx import Presentation
    except ImportError:
        return json.dumps({"error": "python-pptx not installed"})

    prs = Presentation(str(target))
    slides_data: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, 1):
        if i > max_slides:
            break
        text_parts: list[str] = []
        if slide.shapes.title and slide.shapes.title.text:
            text_parts.append(f"Title: {slide.shapes.title.text}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(f"[Table] {row_text}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_parts.append(f"[Notes] {notes}")
        slides_data.append({"slide_num": i, "image_path": None, "text": "\n".join(text_parts)})

    # Try visual rendering: LibreOffice → PDF → pymupdf → PNG
    render_cache = workspace / ".render_cache" / hashlib.md5(str(target).encode()).hexdigest()
    render_cache.mkdir(parents=True, exist_ok=True)
    try:
        import fitz  # pymupdf  # noqa: F401

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            for sd in slides_data:
                sd["image_path"] = "[Install LibreOffice for visual rendering]"
            return json.dumps(slides_data, indent=2)

        pdf_result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(render_cache),
                str(target),
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        pdf_path = render_cache / (target.stem + ".pdf")
        if pdf_path.exists():
            doc = fitz.open(str(pdf_path))
            for idx in range(min(len(doc), max_slides)):
                page = doc[idx]
                pix = page.get_pixmap(dpi=150)
                png_path = render_cache / f"slide_{idx + 1}.png"
                pix.save(str(png_path))
                if idx < len(slides_data):
                    slides_data[idx]["image_path"] = str(png_path)
            doc.close()
        else:
            for sd in slides_data:
                sd["image_path"] = f"[PDF conversion failed: {pdf_result.stderr[:200]}]"
    except ImportError:
        for sd in slides_data:
            sd["image_path"] = "[Install pymupdf for visual rendering]"
    except subprocess.TimeoutExpired:
        for sd in slides_data:
            sd["image_path"] = "[LibreOffice conversion timed out]"
    except Exception as e:
        for sd in slides_data:
            sd["image_path"] = f"[Render error: {e}]"

    return json.dumps(slides_data, indent=2)


def _exec_render_pdf_pages(workspace: Path, args: dict) -> str:
    path = args.get("path", "")
    target = _safe_resolve(workspace, path)
    if target is None:
        return json.dumps({"error": "path outside workspace"})
    if not target.exists():
        return json.dumps({"error": f"File not found: {path}"})
    max_pages = int(args.get("max_pages", 10))

    render_cache = workspace / ".render_cache" / hashlib.md5(str(target).encode()).hexdigest()
    render_cache.mkdir(parents=True, exist_ok=True)
    pages_data: list[dict[str, Any]] = []

    try:
        import fitz

        doc = fitz.open(str(target))
        for i in range(min(len(doc), max_pages)):
            page = doc[i]
            text = page.get_text().strip()
            pix = page.get_pixmap(dpi=150)
            png_path = render_cache / f"page_{i + 1}.png"
            pix.save(str(png_path))
            pages_data.append(
                {
                    "page_num": i + 1,
                    "image_path": str(png_path),
                    "text": text[:5000],
                }
            )
        doc.close()
    except ImportError:
        try:
            import pdfplumber

            with pdfplumber.open(str(target)) as pdf:
                for i, page in enumerate(pdf.pages[:max_pages]):
                    text = page.extract_text() or ""
                    pages_data.append(
                        {
                            "page_num": i + 1,
                            "image_path": "[Install pymupdf for visual rendering]",
                            "text": text[:5000],
                        }
                    )
        except ImportError:
            return json.dumps({"error": "Neither pymupdf nor pdfplumber installed."})

    return json.dumps(pages_data, indent=2)


_OFFICE_EXECUTORS: dict[str, Any] = {
    "read_xlsx_cell": _exec_read_xlsx_cell,
    "read_xlsx_range": _exec_read_xlsx_range,
    "read_docx_content": _exec_read_docx_content,
    "read_pptx_content": _exec_read_pptx_content,
    "render_pptx_slides": _exec_render_pptx_slides,
    "render_pdf_pages": _exec_render_pdf_pages,
}


# ===================================================================
# Tool pack registry
# ===================================================================

_TOOL_PACKS: dict[str, tuple[list[dict], dict[str, Any]]] = {
    "office": (_OFFICE_TOOL_DEFS, _OFFICE_EXECUTORS),
}


def get_judge_tools(
    tool_packs: list[str] | None = None,
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    """Build tool definitions and executor map for the judge.

    Args:
        tool_packs: List of tool pack names to include (e.g. ["office"]).

    Returns:
        (tool_definitions, executor_map) — ready for the agentic loop.
    """
    defs = list(_CORE_TOOL_DEFS)
    executors = dict(_CORE_EXECUTORS)

    for pack_name in tool_packs or []:
        pack = _TOOL_PACKS.get(pack_name)
        if pack is None:
            logger.warning(
                f"Unknown judge tool pack: {pack_name!r} (available: {list(_TOOL_PACKS)})"
            )
            continue
        pack_defs, pack_executors = pack
        defs.extend(pack_defs)
        executors.update(pack_executors)

    return defs, executors


# ===================================================================
# JSON parsing helper
# ===================================================================


def _parse_json_response(text: str) -> dict:
    """Parse JSON from LLM response, handling markdown code blocks."""
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        match = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
        if match:
            return json.loads(match.group(1))
        raise ValueError("Could not parse JSON from judge response")


# ===================================================================
# Default system prompt — used only when the pack does not specify one
# ===================================================================

_DEFAULT_SYSTEM_PROMPT = (
    "You are a thorough grading judge evaluating an AI agent's work product. "
    "You have tools to inspect the agent's workspace files. Use them to examine "
    "the actual deliverables before scoring. Cite specific evidence from the "
    "files in your reasoning. When done, call submit_grade."
)


# ===================================================================
# LLMJudge
# ===================================================================


class LLMJudge:
    """LLM-based grading judge with optional agentic file-reading.

    In agentic mode the judge gets workspace tools (filesystem, shell, search)
    plus any tool packs requested by the pack's grading.yaml.  The judge
    persona/rubric come from the pack — the harness supplies only the mechanism.
    """

    def __init__(self, model_config: ModelConfig):
        self.client = LLMClient(model_config)

    def grade(
        self,
        messages: list[Message],
        rubric: str,
        output_schema: dict[str, Any],
        task_description: str = "",
        workspace_dir: Path | None = None,
        agentic: bool = False,
        system_prompt: str | None = None,
        tool_packs: list[str] | None = None,
    ) -> tuple[float, str]:
        """
        Grade trajectory using LLM judge.

        Args:
            messages: Conversation transcript
            rubric: Grading rubric (from pack's grading.yaml)
            output_schema: Expected JSON schema for judge output
            task_description: Optional task description for context
            workspace_dir: Agent workspace directory (for agentic file reading)
            agentic: If True, judge can call tools to read files before scoring
            system_prompt: Judge persona/instructions (from pack's grading.yaml)
            tool_packs: Tool pack names to load (e.g. ["office"])

        Returns:
            (score 0-1, reasons)
        """
        if agentic and workspace_dir and workspace_dir.exists():
            return self._grade_agentic(
                messages,
                rubric,
                output_schema,
                task_description,
                workspace_dir,
                system_prompt,
                tool_packs,
            )
        return self._grade_single_call(messages, rubric, output_schema, task_description)

    # ------------------------------------------------------------------
    # Single-call grading (original behavior, backward compatible)
    # ------------------------------------------------------------------

    def _grade_single_call(
        self,
        messages: list[Message],
        rubric: str,
        output_schema: dict[str, Any],
        task_description: str,
    ) -> tuple[float, str]:
        transcript_lines = []
        for msg in messages:
            role = msg.role.value.upper() if hasattr(msg.role, "value") else str(msg.role).upper()
            content = msg.content[:500]
            transcript_lines.append(f"{role}: {content}")
            if msg.tool_calls:
                for tc in msg.tool_calls:
                    transcript_lines.append(f"  -> Tool: {tc.name}({tc.arguments})")

        transcript = "\n".join(transcript_lines)

        judge_prompt = (
            f"You are grading an AI agent's performance on a task.\n\n"
            f"Task: {task_description}\n\n"
            f"Rubric:\n{rubric}\n\n"
            f"Conversation Transcript:\n{transcript}\n\n"
            f"Please evaluate the agent's performance according to the rubric and "
            f"provide a score between 0 and 1, along with your reasoning.\n\n"
            f"Respond in JSON format matching this schema:\n"
            f"{json.dumps(output_schema, indent=2)}"
        )

        try:
            result = self.client.generate(
                system=(
                    "You are a fair and thorough grading judge. "
                    "Provide objective assessments based on the rubric."
                ),
                messages=[Message(role="user", content=judge_prompt)],
                temperature=0.0,
            )
            judge_output = _parse_json_response(result.text)
            validate(instance=judge_output, schema=output_schema)
            score = max(0.0, min(1.0, float(judge_output.get("score", 0.0))))
            reasons = judge_output.get("reasoning", judge_output.get("reasons", ""))
            return score, reasons
        except Exception as e:
            return 0.5, f"Judge failed: {e}"

    # ------------------------------------------------------------------
    # Agentic grading (tool-using loop)
    # ------------------------------------------------------------------

    def _grade_agentic(
        self,
        messages: list[Message],
        rubric: str,
        output_schema: dict[str, Any],
        task_description: str,
        workspace_dir: Path,
        system_prompt: str | None = None,
        tool_packs: list[str] | None = None,
    ) -> tuple[float, str]:
        # Build tools from core + requested packs
        tool_defs, executors = get_judge_tools(tool_packs)

        # Build transcript summary
        transcript_lines = []
        for msg in messages:
            role = msg.role.value.upper() if hasattr(msg.role, "value") else str(msg.role).upper()
            content = msg.content[:2000]
            transcript_lines.append(f"{role}: {content}")
            if msg.tool_calls:
                for tc in msg.tool_calls:
                    transcript_lines.append(f"  -> Tool: {tc.name}({tc.arguments})")
        transcript = "\n".join(transcript_lines)

        system_msg = system_prompt or _DEFAULT_SYSTEM_PROMPT

        user_msg = (
            f"Task: {task_description}\n\n"
            f"Rubric:\n{rubric}\n\n"
            f"Agent Conversation Summary:\n{transcript}\n\n"
            f"The agent's workspace is available for inspection. "
            f"Start by listing files, then open and evaluate each deliverable. "
            f"When you have enough evidence, call submit_grade with your score "
            f"(0.0 to 1.0) and detailed reasoning."
        )

        conv: list[Message] = [Message(role="user", content=user_msg)]
        max_turns = 30

        for turn in range(max_turns):
            try:
                result = self.client.generate(
                    system=system_msg,
                    messages=conv,
                    temperature=0.0,
                    tools=tool_defs,
                )
            except Exception as e:
                logger.error(f"Agentic judge LLM call failed: {e}")
                return 0.5, f"Agentic judge failed on turn {turn}: {e}"

            if result.tool_calls:
                conv.append(
                    Message(
                        role="assistant",
                        content=result.text or "",
                        tool_calls=result.tool_calls,
                    )
                )

                for tc in result.tool_calls:
                    tool_name = tc.name
                    try:
                        tool_args = (
                            json.loads(tc.arguments)
                            if isinstance(tc.arguments, str)
                            else tc.arguments
                        )
                    except json.JSONDecodeError:
                        tool_args = {}

                    # Handle submit_grade
                    if tool_name == "submit_grade":
                        score = max(0.0, min(1.0, float(tool_args.get("score", 0.0))))
                        reasoning = tool_args.get("reasoning", "")
                        logger.info(
                            f"Agentic judge submitted grade: {score:.2f} " f"after {turn + 1} turns"
                        )
                        return score, reasoning

                    # Execute tool
                    executor = executors.get(tool_name)
                    if executor:
                        tool_result = executor(workspace_dir, tool_args)
                    else:
                        tool_result = f"Unknown tool: {tool_name}"

                    conv.append(
                        Message(
                            role="tool",
                            content=tool_result,
                            tool_call_id=getattr(tc, "id", tc.name),
                        )
                    )
            else:
                # No tool calls — try to parse as final JSON grade
                if result.text:
                    try:
                        judge_output = _parse_json_response(result.text)
                        validate(instance=judge_output, schema=output_schema)
                        score = max(0.0, min(1.0, float(judge_output.get("score", 0.0))))
                        reasoning = judge_output.get("reasoning", judge_output.get("reasons", ""))
                        return score, reasoning
                    except (json.JSONDecodeError, ValueError):
                        conv.append(Message(role="assistant", content=result.text))
                        conv.append(
                            Message(
                                role="user",
                                content="Please call submit_grade with your final score and reasoning.",
                            )
                        )

        logger.warning("Agentic judge exhausted max turns without submitting grade")
        return 0.5, "Judge did not submit a grade within the turn limit"
